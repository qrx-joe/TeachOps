# -*- coding: utf-8 -*-
"""生成 TeachOps GOAI 初赛 10 页方案 PPTX。

用法：python scripts/make_ppt.py
输出：submission/TeachOps_GOAI_初赛方案.pptx
内容来源：submission/PPT逐页文案.md（步骤 6 产物）。
live 截图占位框在步骤 8/9/10 运行后替换为真实截图，并更新证据索引。
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def _set_dash(shp):
    """占位框虚线边框（本机 python-pptx 无 line 枚举，直接写 XML）。"""
    ln = shp.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))

# ---------- 主题 ----------
FONT = "微软雅黑"
PRIMARY = RGBColor(0x1F, 0x4E, 0x79)   # 深蓝：标题、主结构
DARK = RGBColor(0x22, 0x2E, 0x38)      # 正文
GRAY = RGBColor(0x5A, 0x6B, 0x7B)      # 次要说明
ACCENT = RGBColor(0xC5, 0x5A, 0x11)    # 橙：人工节点 / 占位 / 强调
ACCENT_BG = RGBColor(0xFD, 0xF2, 0xE9)
LIGHT = RGBColor(0xEA, 0xF1, 0xF8)     # 浅蓝底
CARD = RGBColor(0xF2, 0xF6, 0xFA)      # 卡片底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "submission", "TeachOps_GOAI_初赛方案.pptx")


def _run(p, text, size, bold=False, color=DARK):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.name = FONT
    f.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)
    return r


def text_block(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, bold, color, align), ...] 每项一段。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        _run(p, text, size, bold, color)
    return box


def card(slide, x, y, w, h, fill=CARD, line_color=None, dash=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
        if dash:
            _set_dash(shp)
    shp.shadow.inherit = False
    shp.text_frame.word_wrap = True
    return shp


def flow_box(slide, x, y, w, h, text, fill=PRIMARY, size=12):
    shp = card(slide, x, y, w, h, fill=fill)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, text, size, True, WHITE)
    return shp


def arrow(slide, x, y):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.2), Inches(0.26))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xB8, 0xC6, 0xD3)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(slide, title, conclusion, tag="design"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.42), Inches(0.12), Inches(0.52))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False
    text_block(slide, 0.80, 0.34, 10.4, 0.7, [(title, 23, True, DARK, PP_ALIGN.LEFT)])
    text_block(slide, 0.80, 1.02, 11.6, 0.75, [("结论：" + conclusion, 14.5, True, PRIMARY, PP_ALIGN.LEFT)])
    chip = card(slide, 11.35, 0.40, 1.45, 0.4, fill=LIGHT)
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, tag, 10.5, True, PRIMARY)


def footer(slide, n):
    text_block(slide, 0.55, 7.14, 8.0, 0.3, [("研序 TeachOps ｜ GOAI 2026 · Agent Infra 赛道初赛", 9, False, GRAY, PP_ALIGN.LEFT)])
    text_block(slide, 11.8, 7.14, 1.0, 0.3, [(f"{n} / 10", 9, False, GRAY, PP_ALIGN.RIGHT)])


def placeholder(slide, x, y, w, h, label):
    shp = card(slide, x, y, w, h, fill=RGBColor(0xFA, 0xFB, 0xFC), line_color=ACCENT, dash=True)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, label, 12.5, True, ACCENT)
    return shp


def style_table(tbl, widths, header_cells, body_rows, body_size=11.5):
    for i, wd in enumerate(widths):
        tbl.columns[i].width = Inches(wd)
    def fill_cell(cell, text, size, bold, color, bg):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _run(p, text, size, bold, color)
    for c, text in enumerate(header_cells):
        fill_cell(tbl.cell(0, c), text, 12.5, True, WHITE, PRIMARY)
    for r, row in enumerate(body_rows, start=1):
        bg = WHITE if r % 2 == 1 else CARD
        for c, text in enumerate(row):
            fill_cell(tbl.cell(r, c), text, body_size, c == 0, DARK if c else PRIMARY, bg)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------- P1 项目定义（封面式） ----------
    s = new_slide(prs)
    chip = card(s, 0.55, 0.5, 4.1, 0.44, fill=PRIMARY)
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "GOAI 2026 · Agent Infra 赛道 ｜ 初赛方案", 11, True, WHITE)
    text_block(s, 0.55, 1.25, 12.2, 1.1, [("研序 TeachOps：多 Agent 教学设计预审与质量治理系统", 27, True, DARK, PP_ALIGN.LEFT)])
    text_block(s, 0.55, 2.35, 12.2, 0.8, [("教学设计预审从“凭经验改”变成“证据可回溯、缺据即阻断”的多 Agent 流程。", 15.5, True, PRIMARY, PP_ALIGN.LEFT)])
    cards = [
        ("用户", "师范院校 / 教师培训导师"),
        ("任务", "教学设计预审：提交初稿 → 取证 → 修订 → 稽核 → 批准"),
        ("作用", "每条修改建议可回溯证据；缺关键证据即阻断流程"),
    ]
    for i, (t, d) in enumerate(cards):
        x = 0.55 + i * 4.13
        c = card(s, x, 3.35, 3.93, 1.5)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.14)
        p = tf.paragraphs[0]
        _run(p, t, 13, True, PRIMARY)
        p2 = tf.add_paragraph()
        _run(p2, d, 12, False, DARK)
    text_block(s, 0.55, 5.3, 12.2, 1.2, [
        ("问题：初稿常缺课标依据、目标与评价脱节、学情回应不足，人工预审重复且口径不一。", 12, False, GRAY, PP_ALIGN.LEFT),
    ])
    chip2 = card(s, 11.35, 6.55, 1.45, 0.4, fill=LIGHT)
    tf = chip2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "design", 10.5, True, PRIMARY)
    text_block(s, 0.55, 7.14, 8.0, 0.3, [("研序 TeachOps ｜ GOAI 2026 · Agent Infra 赛道初赛", 9, False, GRAY, PP_ALIGN.LEFT)])
    text_block(s, 11.8, 7.14, 1.0, 0.3, [("1 / 10", 9, False, GRAY, PP_ALIGN.RIGHT)])

    # ---------- P2 场景 ----------
    s = new_slide(prs)
    header(s, "教学设计预审场景", "导师在 Word、聊天和通用大模型之间来回，四类问题反复出现。")
    pains = [
        ("痛点 1", "导师反复检查相同的基础缺陷"),
        ("痛点 2", "AI 建议没有可核对的课程依据"),
        ("痛点 3", "修改前后差异与责任记录难以追踪"),
        ("痛点 4", "证据不足时，通用工具仍生成“看似完整”的答案"),
    ]
    for i, (t, d) in enumerate(pains):
        x = 0.55 + (i % 2) * 6.25
        y = 2.0 + (i // 2) * 1.45
        c = card(s, x, y, 6.0, 1.25)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        _run(p, t, 12.5, True, ACCENT)
        p2 = tf.add_paragraph()
        _run(p2, d, 13, False, DARK)
    bar = card(s, 0.55, 5.1, 12.2, 0.8, fill=ACCENT_BG)
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    _run(p, "缺口：通用大模型不维护证据链、不做独立审计、无人工确认边界。", 13.5, True, ACCENT)
    footer(s, 2)

    # ---------- P3 固定课例和输入 ----------
    s = new_slide(prs)
    header(s, "固定课例和输入", "一个三年级数学课例、四件套输入，全部材料可公开。")
    c = card(s, 0.55, 2.0, 3.7, 3.0, fill=LIGHT)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    for i, (t, sz, b, col) in enumerate([
        ("固定课例", 13, True, PRIMARY),
        ("小学三年级数学", 13, False, DARK),
        ("《分数的初步认识》", 14, True, DARK),
        ("第 1 课时 · 40 分钟", 12, False, DARK),
        ("两条演示链路：正常 / 缺证据", 12, False, GRAY),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        _run(p, t, sz, b, col)
    inputs = [
        ("课程标准来源", "教育部 2022 年版课标短引用 + 官方链接（教材〔2022〕2 号）"),
        ("初版教学设计", "团队原创，含 3 个植入的可审计缺陷"),
        ("聚合学情摘要", "synthetic: true —— 无任何学生个人信息"),
        ("规则包", "5 条检查规则，独立于课程标准本身"),
    ]
    for i, (t, d) in enumerate(inputs):
        y = 2.0 + i * 0.78
        c = card(s, 4.45, y, 8.3, 0.66)
        tf = c.text_frame
        tf.margin_left = Inches(0.18)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _run(p, t + "：", 12.5, True, PRIMARY)
        _run(p, d, 12, False, DARK)
    text_block(s, 0.55, 5.3, 12.2, 0.9, [
        ("合规：全部材料可公开；课标仅短引用并附官方链接，不分发全文；学情数据为合成聚合数据。", 12, False, GRAY, PP_ALIGN.LEFT),
    ])
    footer(s, 3)

    # ---------- P4 端到端闭环 ----------
    s = new_slide(prs)
    header(s, "端到端闭环", "证据 → 修订 → 稽核 → 人工批准，一条流水线四道关卡。")
    steps = [
        ("提交课例（导师）", PRIMARY),
        ("Evidence 建包\nREADY / BLOCKED", PRIMARY),
        ("Design 修订\n引用 evidence_id", PRIMARY),
        ("Audit 稽核\n逐规则判定", PRIMARY),
        ("导师批准 / 驳回", ACCENT),
    ]
    for i, (t, col) in enumerate(steps):
        x = 0.55 + i * 2.42
        shp = card(s, x, 2.0, 2.2, 0.95, fill=col)
        tf = shp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.05)
        for j, line in enumerate(t.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            _run(p, line, 11.5 if j == 0 else 10, j == 0, WHITE)
        if i < 4:
            arrow(s, x + 2.24, 2.33)
    arts = ["evidence_packet.json", "revision.md", "audit_report.json", "review_decision.md"]
    for i, a in enumerate(arts):
        x = 0.55 + i * 3.1
        c = card(s, x, 3.2, 2.9, 0.5, fill=CARD)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, a, 11, False, DARK)
    text_block(s, 0.55, 3.95, 12.2, 0.5, [
        ("人工边界：正式通过只能由导师决定，Agent 无发布权；BLOCKED 时流程停止并给出补证清单。", 12.5, True, DARK, PP_ALIGN.LEFT),
    ])
    placeholder(s, 0.55, 4.6, 12.2, 2.1, "【live 截图占位】AgentTeams Team Room 协作记录 —— 步骤 9 运行后替换（当前为 design 流程图）")
    footer(s, 4)

    # ---------- P5 四个 Agent Identity ----------
    s = new_slide(prs)
    header(s, "四个 Agent Identity", "四个角色职责互斥，禁止动作写进 Identity，可逐条检查。")
    tbl = s.shapes.add_table(5, 4, Inches(0.55), Inches(1.95), Inches(12.2), Inches(3.6)).table
    style_table(
        tbl,
        [2.1, 3.5, 3.0, 3.6],
        ["角色", "职责", "产出", "关键禁止"],
        [
            ["Manager", "拆解任务、调度、传递文件引用", "任务计划、状态摘要", "不判断教学质量、不自行批准"],
            ["Evidence Agent", "建证据包、列缺失项", "evidence_packet.json", "不写教案、不补写来源"],
            ["Design Agent", "依据证据修订", "revision.md", "不引用包外依据、不发布版本"],
            ["Audit Agent", "独立逐规则稽核", "audit_report.json", "不修改教案、不代用户批准"],
        ],
    )
    text_block(s, 0.55, 5.8, 12.2, 0.9, [
        ("每个 Identity 含 purpose / inputs / outputs / allowed_actions / forbidden_actions / handoff_to / failure_behavior（agents/*.md）；人类导师是审批者，不是第五个 Agent。", 12, False, GRAY, PP_ALIGN.LEFT),
    ])
    footer(s, 5)

    # ---------- P6 三个 Skill ----------
    s = new_slide(prs)
    header(s, "三个 Skill", "每个 Skill 有契约、有失败样例、可脱离 Agent 单独测试。")
    tbl = s.shapes.add_table(4, 3, Inches(0.55), Inches(1.95), Inches(12.2), Inches(2.9)).table
    style_table(
        tbl,
        [3.5, 4.5, 4.2],
        ["skill_id", "输入 → 输出", "失败契约要点"],
        [
            ["build-evidence-packet", "课标来源 + 学情 → evidence_packet.json", "关键证据缺失即 BLOCKED，不建 READY 包"],
            ["revise-lesson-with-evidence", "初稿 + READY 包 + 规则 → revision.md", "包 BLOCKED 拒绝执行；无据建议入“待补证”区"],
            ["audit-lesson-alignment", "候选设计 + 包 + 规则 → audit_report.json", "判定绑定 rule_id + evidence_id；输入不全不给结论"],
        ],
    )
    text_block(s, 0.55, 5.15, 12.2, 1.3, [
        ("契约含输入输出 JSON Schema、调用条件、权限、超时、错误码、安全说明与正反样例（skills/*/contract.md）。", 12, False, GRAY, PP_ALIGN.LEFT),
        ("可复用性：三者只约束“证据引用纪律 + 审计纪律”，与学科无关，可复用于任意教研文档流水线。", 12, False, GRAY, PP_ALIGN.LEFT),
    ])
    footer(s, 6)

    # ---------- P7 AgentTeams 映射 ----------
    s = new_slide(prs)
    header(s, "AgentTeams 映射", "Manager + 三个 Worker 落在 AgentTeams Team Room，上下文只传文件引用和短状态。")
    tbl = s.shapes.add_table(7, 2, Inches(0.55), Inches(1.9), Inches(12.2), Inches(3.3)).table
    style_table(
        tbl,
        [4.3, 7.9],
        ["TeachOps 概念", "AgentTeams 映射"],
        [
            ["Manager", "Team Leader / Manager"],
            ["Evidence / Design / Audit", "三个 Worker（各挂载对应 Skill 契约）"],
            ["教研任务", "Team Room 中的固定任务"],
            ["上下文传递", "共享文件 + 文件引用（不在消息中复制长文）"],
            ["状态追踪", "房间消息、Agent 状态、产物状态字段"],
            ["人工审批", "用户在 Team Room 中批准 / 驳回"],
        ],
        body_size=11,
    )
    placeholder(s, 0.55, 5.45, 12.2, 1.35, "【live 截图占位】团队配置与 Team Room 成员 —— 步骤 8 完成后替换")
    footer(s, 7)

    # ---------- P8 缺证据异常分支 ----------
    s = new_slide(prs)
    header(s, "缺证据异常分支", "缺关键课标证据 → BLOCKED → 不生成修订 → 导师收到补证清单。")
    steps8 = [
        ("① 异常样例 missing-evidence-case", "与正常样例唯一差异：删除 curriculum-source.md（关键课标证据）"),
        ("② Evidence Agent 返回 BLOCKED", "status: BLOCKED + missing_items，不产出 READY 证据包"),
        ("③ Manager 停止调用 Design Agent", "房间展示停止原因与补证要求；不用模型记忆补写来源，不把部分执行显示成成功"),
    ]
    for i, (t, d) in enumerate(steps8):
        y = 1.95 + i * 1.05
        c = card(s, 0.55, y, 6.9, 0.9)
        tf = c.text_frame
        tf.margin_left = Inches(0.16)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        _run(p, t, 12.5, True, PRIMARY)
        p2 = tf.add_paragraph()
        _run(p2, d, 11.5, False, DARK)
    c = card(s, 7.65, 1.95, 5.1, 2.95, fill=LIGHT)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.16)
    for i, (t, sz, b, col) in enumerate([
        ("fixture replay 回放", 13, True, PRIMARY),
        ("demo/missing-evidence-case/expected-output/", 11, False, DARK),
        ("evidence_packet.json", 11, False, DARK),
        ("内容：BLOCKED + 缺失项清单 + next_action", 11.5, False, DARK),
        ("约定：不制作仿真的 AgentTeams 聊天截图", 11.5, False, ACCENT),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        _run(p, t, sz, b, col)
    placeholder(s, 0.55, 5.25, 12.2, 1.5, "【live 截图占位】BLOCKED 运行记录 —— 步骤 10 运行后替换；未跑通时引用 fixture 并如实标注")
    footer(s, 8)

    # ---------- P9 开放与安全边界 ----------
    s = new_slide(prs)
    header(s, "开放与安全边界", "MIT 开源四类资产，四条安全边界贯穿全部材料。")
    c = card(s, 0.55, 2.0, 5.9, 4.0)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    _run(p, "开源范围（MIT）", 14, True, PRIMARY)
    for t in [
        "Agent Identity（四个角色定义）",
        "Skill 契约（含输入输出 JSON Schema）",
        "固定课例样例与期望输出（正常 + 缺证据）",
        "README、运行证据索引与运行手册",
        "PPT 生成脚本（scripts/make_ppt.py，可复现）",
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(7)
        _run(p, "· " + t, 12.5, False, DARK)
    c = card(s, 6.65, 2.0, 6.1, 4.0)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    _run(p, "四条安全边界", 14, True, PRIMARY)
    for t in [
        "① API Key 只存本机，不进聊天、仓库、截图与 PPT",
        "② 学情为合成聚合数据，无学生个人信息",
        "③ 课标仅短引用 + 官方链接，不分发全文",
        "④ 批准 / 驳回权在人类导师，Agent 不能自行结束流程",
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(9)
        _run(p, t, 12.5, False, DARK)
    footer(s, 9)

    # ---------- P10 当前进展、限制和复赛计划 ----------
    s = new_slide(prs)
    header(s, "当前进展、限制和复赛计划", "初赛交付材料齐备，未做的事项明确列出，复赛路线可检验。")
    c = card(s, 0.55, 1.9, 12.2, 0.95, fill=LIGHT)
    tf = c.text_frame
    tf.margin_left = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _run(p, "已完成：", 13, True, PRIMARY)
    _run(p, "四个 Agent Identity ｜ 三个 Skill 契约 ｜ 两套固定样例与期望输出 ｜ 500 字作品简介 ｜ 本 PPT（生成脚本可复现）", 12.5, False, DARK)
    c = card(s, 0.55, 3.0, 12.2, 0.8, fill=ACCENT_BG)
    tf = c.text_frame
    tf.margin_left = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _run(p, "AgentTeams 运行状态：", 12.5, True, ACCENT)
    _run(p, "烟雾测试与两条链路运行进行中，live 证据以 docs/运行证据索引.md 登记为准；未跑通部分以 fixture replay / design 如实呈现。", 12, False, DARK)
    c = card(s, 0.55, 4.0, 5.9, 2.4)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    _run(p, "限制", 13.5, True, PRIMARY)
    for t in [
        "无真实导师 / 学生数据，效果未在真实组织验证",
        "单一课例（三年级数学），未做多学科扩展",
        "AgentTeams 未在本机长期稳定性验证",
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(6)
        _run(p, "· " + t, 12, False, DARK)
    c = card(s, 6.65, 4.0, 6.1, 2.4)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    _run(p, "复赛计划", 13.5, True, PRIMARY)
    for t in [
        "教研任务 Web 工作台",
        "FastAPI + PostgreSQL 正式状态机",
        "版本 Diff 与回滚、审批与审计事件",
        "知识库 RAG / MCP Adapter",
        "1-3 名真实导师的反馈验证",
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        _run(p, "· " + t, 12, False, DARK)
    footer(s, 10)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("saved:", OUT)


if __name__ == "__main__":
    build()
